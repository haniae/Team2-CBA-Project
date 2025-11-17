"""PowerPoint presentation builder for portfolio reports.

Generates CFI-style 12-slide professional presentations with portfolio analysis,
optimization results, exposure analysis, attribution, and scenario stress tests.
"""

from __future__ import annotations

import io
import logging
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PowerPoint export will be disabled.")

from .portfolio import (
    OptimizationResult,
    AttributionResult,
    ExposureAnalysis,
    ScenarioResult,
)

LOGGER = logging.getLogger(__name__)

# Theme colors
NAVY_DEEP = RGBColor(20, 40, 80) if PPTX_AVAILABLE else None
NAVY_LIGHT = RGBColor(70, 130, 180) if PPTX_AVAILABLE else None
BLUE_LIGHT = RGBColor(173, 216, 230) if PPTX_AVAILABLE else None
WHITE = RGBColor(255, 255, 255) if PPTX_AVAILABLE else None
GREY_LIGHT = RGBColor(200, 200, 200) if PPTX_AVAILABLE else None
GREY_DARK = RGBColor(80, 80, 80) if PPTX_AVAILABLE else None


def build_portfolio_powerpoint(
    database_path: Path,
    portfolio_id: str,
    optimization_result: Optional[OptimizationResult] = None,
    attribution_result: Optional[AttributionResult] = None,
    exposure_result: Optional[ExposureAnalysis] = None,
    scenario_results: Optional[List[ScenarioResult]] = None,
) -> bytes:
    """Build a 12-slide PowerPoint presentation for a portfolio.
    
    Args:
        database_path: Path to the database
        portfolio_id: Portfolio identifier
        optimization_result: Optional optimization results
        attribution_result: Optional attribution analysis results
        exposure_result: Optional exposure analysis results
        scenario_results: Optional scenario stress test results
        
    Returns:
        bytes: PowerPoint file content
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required for PowerPoint export. Install with: pip install python-pptx")
    
    # Fetch portfolio data
    from . import database
    metadata = database.fetch_portfolio_metadata(database_path, portfolio_id)
    if not metadata:
        raise ValueError(f"Portfolio {portfolio_id} not found")
    
    holdings = database.fetch_portfolio_holdings(database_path, portfolio_id)
    portfolio_name = metadata.name if hasattr(metadata, 'name') else portfolio_id
    created_at = metadata.created_at if hasattr(metadata, 'created_at') else datetime.now()
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Cover
    _add_cover_slide(prs, portfolio_name, created_at)
    
    # SLIDE 2: Executive Summary
    _add_executive_summary_slide(prs, metadata, holdings)
    
    # SLIDE 3: Portfolio Holdings
    _add_holdings_slide(prs, holdings)
    
    # SLIDE 4: Sector Exposure
    if exposure_result:
        _add_sector_exposure_slide(prs, exposure_result)
    else:
        _add_blank_slide(prs, "Sector Exposure", "Run exposure analysis to view sector breakdown.")
    
    # SLIDE 5: Factor Exposure
    if exposure_result:
        _add_factor_exposure_slide(prs, exposure_result)
    else:
        _add_blank_slide(prs, "Factor Exposure", "Run exposure analysis to view factor exposures.")
    
    # SLIDE 6: Performance Attribution
    if attribution_result:
        _add_attribution_slide(prs, attribution_result)
    else:
        _add_blank_slide(prs, "Performance Attribution", "Run attribution analysis to view performance breakdown.")
    
    # SLIDE 7: Optimization Results
    if optimization_result:
        _add_optimization_slide(prs, optimization_result)
    else:
        _add_blank_slide(prs, "Portfolio Optimization", "Run optimization to view recommended changes.")
    
    # SLIDE 8: Proposed Trades
    if optimization_result and optimization_result.proposed_trades:
        _add_proposed_trades_slide(prs, optimization_result)
    else:
        _add_blank_slide(prs, "Proposed Trades", "Run optimization to view recommended trades.")
    
    # SLIDE 9: Scenario Analysis
    if scenario_results:
        _add_scenario_analysis_slide(prs, scenario_results)
    else:
        _add_blank_slide(prs, "Scenario Analysis", "Run scenario stress tests to view impact analysis.")
    
    # SLIDE 10: Risk Metrics
    _add_risk_metrics_slide(prs, database_path, portfolio_id)
    
    # SLIDE 11: Concentration Metrics
    if exposure_result:
        _add_concentration_slide(prs, exposure_result)
    else:
        _add_blank_slide(prs, "Concentration Metrics", "Run exposure analysis to view concentration metrics.")
    
    # SLIDE 12: Appendix
    _add_appendix_slide(prs, portfolio_id, created_at)
    
    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _add_cover_slide(prs: Presentation, portfolio_name: str, created_at: Any) -> None:
    """Add cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY_DEEP
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = portfolio_name
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Portfolio Analysis Report"
    p.font.size = Pt(20)
    p.font.color.rgb = BLUE_LIGHT
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    try:
        if isinstance(created_at, datetime):
            date_str = created_at.strftime("%B %d, %Y")
        else:
            dt = datetime.fromisoformat(str(created_at).replace('Z', '+00:00'))
            date_str = dt.strftime("%B %d, %Y")
    except:
        date_str = str(created_at)
    p.text = date_str
    p.font.size = Pt(14)
    p.font.color.rgb = GREY_LIGHT
    p.alignment = PP_ALIGN.CENTER


def _add_executive_summary_slide(prs: Presentation, metadata: Any, holdings: List[Any]) -> None:
    """Add executive summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    # Content
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True
    
    portfolio_name = getattr(metadata, 'name', 'N/A')
    base_currency = getattr(metadata, 'base_currency', 'USD')
    benchmark = getattr(metadata, 'benchmark_index', 'S&P 500')
    
    lines = [
        f"Portfolio Name: {portfolio_name}",
        f"Base Currency: {base_currency}",
        f"Benchmark: {benchmark}",
        f"Total Holdings: {len(holdings)}",
        "",
        "This report provides a comprehensive analysis of portfolio composition,",
        "performance attribution, risk metrics, and optimization opportunities.",
    ]
    
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(14)
        p.space_after = Pt(6)


def _add_holdings_slide(prs: Presentation, holdings: List[Any]) -> None:
    """Add portfolio holdings slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Portfolio Holdings"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    # Table
    table = slide.shapes.add_table(rows=min(len(holdings) + 1, 15), cols=3, 
                                    left=Inches(0.5), top=Inches(1.2), 
                                    width=Inches(9), height=Inches(5.5)).table
    
    # Header
    table.cell(0, 0).text = "Ticker"
    table.cell(0, 1).text = "Weight (%)"
    table.cell(0, 2).text = "Market Value"
    
    for i, cell in enumerate(table.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_DEEP
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Data
    for i, holding in enumerate(holdings[:14]):
        row = i + 1
        weight = getattr(holding, 'weight', None) or 0
        if weight and weight > 1:
            weight = weight / 100.0
        weight_pct = (weight * 100) if weight else 0
        market_value = getattr(holding, 'market_value', None) or 0
        
        table.cell(row, 0).text = getattr(holding, 'ticker', 'N/A')
        table.cell(row, 1).text = f"{weight_pct:.2f}%"
        table.cell(row, 2).text = f"${market_value:,.2f}" if market_value else "N/A"
        
        for cell in table.rows[row].cells:
            cell.text_frame.paragraphs[0].font.size = Pt(11)


def _add_sector_exposure_slide(prs: Presentation, exposure: ExposureAnalysis) -> None:
    """Add sector exposure slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Sector Exposure"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    # Create pie chart if we have sector data
    if exposure.sector_exposure:
        chart_data = ChartData()
        chart_data.categories = list(exposure.sector_exposure.keys())
        chart_data.add_series("Weight", [v * 100 for v in exposure.sector_exposure.values()])
        
        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(2), Inches(1.5), 
                                       Inches(6), Inches(4.5), chart_data).chart
        chart.has_legend = True
        chart.chart_title.text_frame.text = "Portfolio by Sector"


def _add_factor_exposure_slide(prs: Presentation, exposure: ExposureAnalysis) -> None:
    """Add factor exposure slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Factor Exposure"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    if exposure.factor_exposure:
        sorted_factors = sorted(exposure.factor_exposure.items(), 
                               key=lambda x: abs(x[1]), reverse=True)[:10]
        
        table = slide.shapes.add_table(rows=len(sorted_factors) + 1, cols=2,
                                       left=Inches(1), top=Inches(1.5),
                                       width=Inches(8), height=Inches(5)).table
        
        table.cell(0, 0).text = "Factor"
        table.cell(0, 1).text = "Exposure"
        
        for i, (factor, value) in enumerate(sorted_factors):
            row = i + 1
            table.cell(row, 0).text = factor.replace("_", " ").title()
            table.cell(row, 1).text = f"{value:+.3f}"


def _add_attribution_slide(prs: Presentation, attribution: AttributionResult) -> None:
    """Add performance attribution slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Performance Attribution"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    total_alloc = sum(attribution.allocation_effect.values())
    total_select = sum(attribution.selection_effect.values())
    total_interact = sum(attribution.interaction_effect.values())
    
    lines = [
        f"Total Active Return: {attribution.total_active_return:.2%}",
        f"Allocation Effect: {total_alloc:.2%}",
        f"Selection Effect: {total_select:.2%}",
        f"Interaction Effect: {total_interact:.2%}",
    ]
    
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(16)
        p.space_after = Pt(8)


def _add_optimization_slide(prs: Presentation, optimization: OptimizationResult) -> None:
    """Add optimization results slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Optimization Results"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    before = optimization.metrics_before
    after = optimization.metrics_after
    
    lines = [
        "Expected Return:",
        f"  Before: {before.get('expected_return', 0):.2%}",
        f"  After: {after.get('expected_return', 0):.2%}",
        "",
        "Portfolio Variance:",
        f"  Before: {before.get('portfolio_variance', 0):.4f}",
        f"  After: {after.get('portfolio_variance', 0):.4f}",
        "",
        f"Status: {optimization.status}",
    ]
    
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(14)
        p.space_after = Pt(6)


def _add_proposed_trades_slide(prs: Presentation, optimization: OptimizationResult) -> None:
    """Add proposed trades slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Proposed Trades"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    trades = optimization.proposed_trades[:10]  # Top 10
    
    table = slide.shapes.add_table(rows=len(trades) + 1, cols=4,
                                   left=Inches(0.5), top=Inches(1.2),
                                   width=Inches(9), height=Inches(5.5)).table
    
    table.cell(0, 0).text = "Ticker"
    table.cell(0, 1).text = "Action"
    table.cell(0, 2).text = "From (%)"
    table.cell(0, 3).text = "To (%)"
    
    for i, trade in enumerate(trades):
        row = i + 1
        table.cell(row, 0).text = trade['ticker']
        table.cell(row, 1).text = trade['action']
        table.cell(row, 2).text = f"{trade['from_weight']:.2f}%"
        table.cell(row, 3).text = f"{trade['to_weight']:.2f}%"


def _add_scenario_analysis_slide(prs: Presentation, scenarios: List[ScenarioResult]) -> None:
    """Add scenario analysis slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Scenario Analysis"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    for i, scenario in enumerate(scenarios[:5]):  # Top 5 scenarios
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"Scenario {i+1}: Portfolio Return = {scenario.portfolio_return:.2%}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.space_after = Pt(4)


def _add_risk_metrics_slide(prs: Presentation, database_path: Path, portfolio_id: str) -> None:
    """Add risk metrics slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Risk Metrics"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    try:
        from .portfolio import calculate_portfolio_beta, get_portfolio_returns
        beta = calculate_portfolio_beta(database_path, portfolio_id)
        returns = get_portfolio_returns(database_path, portfolio_id, 252)
        import numpy as np
        volatility = returns.std() * (252 ** 0.5) if len(returns) > 0 else 0
        
        lines = [
            f"Portfolio Beta: {beta:.2f}",
            f"Annualized Volatility: {volatility:.2%}",
        ]
        
        for i, line in enumerate(lines):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(16)
            p.space_after = Pt(8)
    except Exception as e:
        LOGGER.warning(f"Failed to calculate risk metrics: {e}")
        p = tf.paragraphs[0]
        p.text = "Risk metrics calculation unavailable."
        p.font.size = Pt(14)


def _add_concentration_slide(prs: Presentation, exposure: ExposureAnalysis) -> None:
    """Add concentration metrics slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Concentration Metrics"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    if exposure.concentration_metrics:
        metrics_list = list(exposure.concentration_metrics.items())[:8]
        table = slide.shapes.add_table(rows=len(metrics_list) + 1, cols=2,
                                       left=Inches(1), top=Inches(1.5),
                                       width=Inches(8), height=Inches(5)).table
        
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        
        for i, (metric, value) in enumerate(metrics_list):
            row = i + 1
            table.cell(row, 0).text = metric.replace("_", " ").title()
            table.cell(row, 1).text = f"{value:.3f}"


def _add_appendix_slide(prs: Presentation, portfolio_id: str, created_at: Any) -> None:
    """Add appendix slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = "Appendix"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    try:
        if isinstance(created_at, datetime):
            date_str = created_at.strftime("%B %d, %Y at %I:%M %p UTC")
        else:
            dt = datetime.fromisoformat(str(created_at).replace('Z', '+00:00'))
            date_str = dt.strftime("%B %d, %Y at %I:%M %p UTC")
    except:
        date_str = str(created_at)
    
    lines = [
        f"Portfolio ID: {portfolio_id}",
        f"Report Generated: {date_str}",
        "",
        "This report was generated by FinanlyzeOS IVPA (Investment Portfolio Agent).",
        "",
        "For questions or support, please contact your portfolio administrator.",
    ]
    
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(12)
        p.space_after = Pt(6)


def _add_blank_slide(prs: Presentation, title_text: str, message: str) -> None:
    """Add a blank slide with title and message."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_DEEP
    
    content = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = message
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER


