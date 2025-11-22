"""Document text extraction for various file formats."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional, Tuple

LOGGER = logging.getLogger(__name__)


def extract_text_from_file(file_path: Path, filename: Optional[str] = None) -> Tuple[Optional[str], Optional[str]]:
    """
    Extract text from various file formats.
    
    Args:
        file_path: Path to the file
        filename: Optional filename (for type detection)
        
    Returns:
        Tuple of (extracted_text, file_type) or (None, None) if extraction fails
    """
    if not file_path.exists():
        LOGGER.warning(f"File not found: {file_path}")
        return None, None
    
    # Determine file type from extension
    file_ext = (file_path.suffix or "").lower()
    if filename:
        filename_ext = Path(filename).suffix.lower()
        if filename_ext:
            file_ext = filename_ext
    
    # Try to extract based on file type
    try:
        if file_ext == '.pdf':
            return extract_text_from_pdf(file_path), 'pdf'
        elif file_ext in ['.docx', '.doc']:
            return extract_text_from_docx(file_path), 'docx'
        elif file_ext in ['.pptx', '.ppt']:
            return extract_text_from_pptx(file_path), 'pptx'
        elif file_ext == '.txt':
            return extract_text_from_txt(file_path), 'text'
        elif file_ext == '.csv':
            return extract_text_from_csv(file_path), 'csv'
        elif file_ext in ['.xlsx', '.xls']:
            return extract_text_from_excel(file_path), 'excel'
        elif file_ext == '.json':
            return extract_text_from_json(file_path), 'json'
        elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
            return extract_text_from_image(file_path), 'image'
        elif file_ext in ['.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.rb', '.go', '.rs', '.php', '.html', '.css', '.xml', '.yaml', '.yml']:
            return extract_text_from_txt(file_path), 'code'
        else:
            # Try as text file as fallback
            try:
                text = extract_text_from_txt(file_path)
                if text and len(text.strip()) > 10:
                    return text, 'text'
            except Exception:
                pass
            return None, 'unknown'
    except Exception as e:
        LOGGER.error(f"Error extracting text from {file_path}: {e}", exc_info=True)
        return None, None


def extract_text_from_pdf(file_path: Path) -> Optional[str]:
    """Extract text from PDF file."""
    try:
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip() if text.strip() else None
        except ImportError:
            # Try pypdf instead
            try:
                import pypdf
                with open(file_path, 'rb') as f:
                    pdf_reader = pypdf.PdfReader(f)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                    return text.strip() if text.strip() else None
            except ImportError:
                # Try pdfplumber
                try:
                    import pdfplumber
                    with pdfplumber.open(file_path) as pdf:
                        text = ""
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                        return text.strip() if text.strip() else None
                except ImportError:
                    LOGGER.warning("PDF extraction libraries not available. Install PyPDF2, pypdf, or pdfplumber.")
                    return None
    except Exception as e:
        LOGGER.error(f"Error extracting text from PDF {file_path}: {e}", exc_info=True)
        return None


def extract_text_from_docx(file_path: Path) -> Optional[str]:
    """Extract text from Word document (.docx)."""
    try:
        try:
            from docx import Document
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += "\n" + cell.text
            return text.strip() if text.strip() else None
        except ImportError:
            LOGGER.warning("python-docx not available. Install with: pip install python-docx")
            return None
    except Exception as e:
        LOGGER.error(f"Error extracting text from DOCX {file_path}: {e}", exc_info=True)
        return None


def extract_text_from_txt(file_path: Path) -> Optional[str]:
    """Extract text from plain text file."""
    try:
        # Try multiple encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1', 'utf-16']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                    text = f.read()
                    if text.strip():
                        return text.strip()
            except (UnicodeDecodeError, UnicodeError):
                continue
        # If all encodings fail, try binary read with error replacement
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
                text = content.decode('utf-8', errors='replace')
                if text.strip():
                    return text.strip()
        except Exception:
            pass
        return None
    except Exception as e:
        LOGGER.error(f"Error extracting text from TXT {file_path}: {e}", exc_info=True)
        return None


def extract_text_from_csv(file_path: Path) -> Optional[str]:
    """Extract text from CSV file."""
    try:
        import csv
        text_lines = []
        with open(file_path, 'r', encoding='utf-8', errors='replace', newline='') as f:
            reader = csv.reader(f)
            for row in reader:
                text_lines.append(", ".join(str(cell) for cell in row))
        text = "\n".join(text_lines)
        return text.strip() if text.strip() else None
    except Exception as e:
        LOGGER.error(f"Error extracting text from CSV {file_path}: {e}", exc_info=True)
        # Fallback to plain text extraction
        return extract_text_from_txt(file_path)


def extract_text_from_excel(file_path: Path) -> Optional[str]:
    """Extract text from Excel file."""
    try:
        try:
            import pandas as pd
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_lines = []
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                text_lines.append(f"Sheet: {sheet_name}")
                text_lines.append(df.to_string())
                text_lines.append("")
            text = "\n".join(text_lines)
            return text.strip() if text.strip() else None
        except ImportError:
            LOGGER.warning("pandas or openpyxl not available. Install with: pip install pandas openpyxl")
            return None
    except Exception as e:
        LOGGER.error(f"Error extracting text from Excel {file_path}: {e}", exc_info=True)
        return None


def extract_text_from_json(file_path: Path) -> Optional[str]:
    """Extract text from JSON file."""
    try:
        import json
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            data = json.load(f)
            # Convert JSON to readable text
            text = json.dumps(data, indent=2, ensure_ascii=False)
            return text.strip() if text.strip() else None
    except Exception as e:
        LOGGER.error(f"Error extracting text from JSON {file_path}: {e}", exc_info=True)
        # Fallback to plain text extraction
        return extract_text_from_txt(file_path)


def extract_text_from_pptx(file_path: Path) -> Optional[str]:
    """Extract text from PowerPoint file (.pptx)."""
    try:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            
            # Extract text from slides
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content.append(f"\n--- Slide {slide_num} ---\n")
                
                # Extract from shapes (text boxes, etc.)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content.append(shape.text)
                
                # Extract from tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells]
                            text_content.append(" | ".join(row_data))
            
            text = "\n".join(text_content)
            return text.strip() if text.strip() else None
        except ImportError:
            LOGGER.warning("python-pptx not available. Install with: pip install python-pptx")
            return None
    except Exception as e:
        LOGGER.error(f"Error extracting text from PPTX {file_path}: {e}", exc_info=True)
        return None


def extract_text_from_image(file_path: Path) -> Optional[str]:
    """Extract text from image file using OCR."""
    try:
        # Try Tesseract OCR first
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            if text and text.strip():
                return f"[IMAGE CONTENT EXTRACTED VIA OCR]\n{text.strip()}\n[/IMAGE CONTENT]"
            else:
                return "[IMAGE FILE] - No text detected via OCR. The image has been stored for reference."
        except ImportError:
            LOGGER.warning("pytesseract or PIL not available. Install with: pip install pytesseract pillow")
            # Fallback: return placeholder
            return "[IMAGE FILE] - OCR not available. The image has been stored for reference. Install pytesseract for text extraction."
        except Exception as e:
            LOGGER.error(f"Error during OCR extraction: {e}", exc_info=True)
            return "[IMAGE FILE] - OCR extraction failed. The image has been stored for reference."
    except Exception as e:
        LOGGER.error(f"Error extracting text from image {file_path}: {e}", exc_info=True)
        return "[IMAGE FILE] - Error processing image. The file has been stored for reference."

