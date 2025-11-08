"""
Template Processing System

Extracts structure from PowerPoint/PDF/Excel templates and fills them with new data.
"""

from __future__ import annotations

import json
import logging
import sqlite3
import uuid
import tempfile
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from . import database
from .document_processor import extract_text_from_file
from .custom_kpis import CustomKPICalculator
from .interactive_modeling import ModelBuilder
from .predictive_analytics import PredictiveAnalytics
from .source_tracer import SourceTracer
import re

LOGGER = logging.getLogger(__name__)


def _json_dumps(payload: Any) -> str:
    """Deterministic JSON serializer for template metadata."""
    return json.dumps(payload, default=str, separators=(",", ":"))


@dataclass
class ReportTemplate:
    """Represents a report template."""
    template_id: str
    user_id: str
    name: str
    file_type: str
    file_path: str
    structure: Dict[str, Any]
    created_at: datetime
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        return {
            "template_id": self.template_id,
            "user_id": self.user_id,
            "name": self.name,
            "file_type": self.file_type,
            "file_path": self.file_path,
            "structure": self.structure,
            "created_at": self.created_at.isoformat(),
        }


@dataclass
class TemplatePlaceholder:
    """Represents a placeholder in a template."""
    placeholder_id: str
    template_id: str
    placeholder_name: str
    data_type: str
    source_metric: Optional[str]
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        return {
            "placeholder_id": self.placeholder_id,
            "template_id": self.template_id,
            "placeholder_name": self.placeholder_name,
            "data_type": self.data_type,
            "source_metric": self.source_metric,
        }


@dataclass
class TemplateRenderResult:
    """Result of filling a template."""

    success: bool
    output_path: Path
    filled_data: Dict[str, Any] = field(default_factory=dict)
    audit_log: List[Dict[str, Any]] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    job_id: Optional[str] = None


class TemplateExtractor:
    """Extracts structure and placeholders from templates."""
    
    def extract_structure(self, file_path: Path, file_type: str) -> Dict[str, Any]:
        """Extract structure from a template file."""
        structure = {
            "file_type": file_type,
            "placeholders": [],
            "sections": [],
        }
        
        if file_type == "pptx":
            structure.update(self._extract_pptx_structure(file_path))
        elif file_type == "pdf":
            structure.update(self._extract_pdf_structure(file_path))
        elif file_type in ["xlsx", "xls"]:
            structure.update(self._extract_excel_structure(file_path))
        else:
            # Fallback: extract text and find placeholders
            text, _ = extract_text_from_file(file_path)
            if text:
                structure["placeholders"] = self._find_placeholders_in_text(text)
        
        return structure
    
    def _extract_pptx_structure(self, file_path: Path) -> Dict[str, Any]:
        """Extract structure from PowerPoint file."""
        structure = {
            "slides": [],
            "placeholders": [],
        }
        
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_info = {
                    "slide_number": slide_idx + 1,
                    "shapes": [],
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        # Look for placeholders
                        placeholders = self._find_placeholders_in_text(text)
                        if placeholders:
                            structure["placeholders"].extend(placeholders)
                        
                        slide_info["shapes"].append({
                            "type": shape.shape_type,
                            "text": text[:200],  # Limit text length
                        })
                
                structure["slides"].append(slide_info)
        except ImportError:
            LOGGER.warning("python-pptx not available for PowerPoint processing")
        except Exception as e:
            LOGGER.error(f"Error extracting PowerPoint structure: {e}", exc_info=True)
        
        return structure
    
    def _extract_pdf_structure(self, file_path: Path) -> Dict[str, Any]:
        """Extract structure from PDF file."""
        structure = {
            "pages": [],
            "placeholders": [],
        }
        
        text, _ = extract_text_from_file(file_path)
        if text:
            # Split by pages (approximate)
            lines = text.split("\n")
            page_size = 50  # Approximate lines per page
            for page_idx in range(0, len(lines), page_size):
                page_text = "\n".join(lines[page_idx:page_idx + page_size])
                placeholders = self._find_placeholders_in_text(page_text)
                if placeholders:
                    structure["placeholders"].extend(placeholders)
                
                structure["pages"].append({
                    "page_number": (page_idx // page_size) + 1,
                    "text_preview": page_text[:500],
                })
        
        return structure
    
    def _extract_excel_structure(self, file_path: Path) -> Dict[str, Any]:
        """Extract structure from Excel file."""
        structure = {
            "worksheets": [],
            "placeholders": [],
        }
        
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_info = {
                    "name": sheet_name,
                    "cells": [],
                }
                
                # Extract cell values and look for placeholders
                for row in sheet.iter_rows(max_row=100, max_col=20):  # Limit to first 100 rows, 20 cols
                    for cell in row:
                        if cell.value:
                            cell_text = str(cell.value)
                            placeholders = self._find_placeholders_in_text(cell_text)
                            if placeholders:
                                structure["placeholders"].extend(placeholders)
                            
                            sheet_info["cells"].append({
                                "cell": cell.coordinate,
                                "value": cell_text[:200],
                            })
                
                structure["worksheets"].append(sheet_info)
        except ImportError:
            LOGGER.warning("openpyxl not available for Excel processing")
        except Exception as e:
            LOGGER.error(f"Error extracting Excel structure: {e}", exc_info=True)
        
        return structure
    
    def _classify_placeholder(self, name: str) -> Tuple[str, str]:
        if ":" in name:
            prefix, remainder = name.split(":", 1)
            prefix_lower = prefix.strip().lower()
            identifier = remainder.strip()
            if prefix_lower == "kpi":
                return "kpi", identifier
            if prefix_lower == "chart":
                return "chart", identifier
            if prefix_lower in {"model", "forecast"}:
                return "model", identifier
            if prefix_lower == "text":
                return "text", identifier
        return "text", name.strip()

    def _find_placeholders_in_text(self, text: str) -> List[Dict[str, str]]:
        """Find placeholders in text (e.g., {{metric_name}}, {ticker}, etc.)."""
        placeholders = []
        
        # Pattern 1: {{placeholder}}
        pattern1 = r'\{\{([^}]+)\}\}'
        matches = re.finditer(pattern1, text)
        for match in matches:
            data_type, identifier = self._classify_placeholder(match.group(1).strip())
            placeholders.append({
                "name": match.group(1).strip(),
                "pattern": match.group(0),
                "type": "double_brace",
                "data_type": data_type,
                "identifier": identifier,
            })
        
        # Pattern 2: {placeholder}
        pattern2 = r'\{([^}]+)\}'
        matches = re.finditer(pattern2, text)
        for match in matches:
            # Skip if already found as double brace
            if match.group(0) not in [p["pattern"] for p in placeholders]:
                data_type, identifier = self._classify_placeholder(match.group(1).strip())
                placeholders.append({
                    "name": match.group(1).strip(),
                    "pattern": match.group(0),
                    "type": "single_brace",
                    "data_type": data_type,
                    "identifier": identifier,
                })
        
        # Pattern 3: [PLACEHOLDER]
        pattern3 = r'\[([A-Z_]+)\]'
        matches = re.finditer(pattern3, text)
        for match in matches:
            data_type, identifier = self._classify_placeholder(match.group(1).strip())
            placeholders.append({
                "name": match.group(1).strip(),
                "pattern": match.group(0),
                "type": "bracket",
                "data_type": data_type,
                "identifier": identifier,
            })
        
        return placeholders


class TemplateFiller:
    """Fills templates with new data."""
    
    def fill_template(
        self,
        template_path: Path,
        template_type: str,
        data: Dict[str, Any],
        output_path: Path
    ) -> bool:
        """Fill a template with data and save to output path."""
        try:
            if template_type == "pptx":
                return self._fill_pptx(template_path, data, output_path)
            elif template_type == "pdf":
                return self._fill_pdf(template_path, data, output_path)
            elif template_type in ["xlsx", "xls"]:
                return self._fill_excel(template_path, data, output_path)
            else:
                return self._fill_text(template_path, data, output_path)
        except Exception as e:
            LOGGER.error(f"Error filling template: {e}", exc_info=True)
            return False
    
    def _fill_pptx(self, template_path: Path, data: Dict[str, Any], output_path: Path) -> bool:
        """Fill PowerPoint template."""
        try:
            from pptx import Presentation
            prs = Presentation(template_path)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape.text = self._replace_placeholders(shape.text, data)
            
            prs.save(output_path)
            return True
        except Exception as e:
            LOGGER.error(f"Error filling PowerPoint: {e}", exc_info=True)
            return False
    
    def _fill_pdf(self, template_path: Path, data: Dict[str, Any], output_path: Path) -> bool:
        """Fill PDF template (creates new PDF with replaced text)."""
        # PDF filling is complex - for now, just copy the file
        # In production, would use a PDF library like reportlab or pdfrw
        try:
            import shutil
            shutil.copy(template_path, output_path)
            LOGGER.warning("PDF filling not fully implemented - file copied as-is")
            return True
        except Exception as e:
            LOGGER.error(f"Error filling PDF: {e}", exc_info=True)
            return False
    
    def _fill_excel(self, template_path: Path, data: Dict[str, Any], output_path: Path) -> bool:
        """Fill Excel template."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(template_path)
            
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            cell.value = self._replace_placeholders(cell.value, data)
            
            wb.save(output_path)
            return True
        except Exception as e:
            LOGGER.error(f"Error filling Excel: {e}", exc_info=True)
            return False
    
    def _fill_text(self, template_path: Path, data: Dict[str, Any], output_path: Path) -> bool:
        """Fill text template."""
        try:
            text = template_path.read_text(encoding='utf-8')
            filled_text = self._replace_placeholders(text, data)
            output_path.write_text(filled_text, encoding='utf-8')
            return True
        except Exception as e:
            LOGGER.error(f"Error filling text: {e}", exc_info=True)
            return False
    
    def _replace_placeholders(self, text: str, data: Dict[str, Any]) -> str:
        """Replace placeholders in text with data values."""
        result = text
        
        # Replace {{placeholder}}
        for key, value in data.items():
            replacement = value
            if isinstance(value, dict):
                replacement = value.get("text") or value.get("image_path") or ""
            result = result.replace(f"{{{{{key}}}}}", str(replacement))
            result = result.replace(f"{{{key}}}", str(replacement))
            result = result.replace(f"[{key.upper()}]", str(replacement))
        
        return result


class ChartGenerator:
    """Generates chart artifacts for template placeholders."""

    def __init__(self, db_path: Path):
        self.db_path = db_path
        self.analytics = PredictiveAnalytics(str(db_path))

    def generate(
        self,
        ticker: str,
        chart_identifier: str,
        context: Dict[str, Any],
    ) -> Tuple[Optional[str], Dict[str, Any], Optional[str]]:
        """Create a chart image for the requested specification.

        Returns:
            Tuple of (file_path, audit_entry, warning)
        """
        try:
            import matplotlib.pyplot as plt  # type: ignore
        except Exception as exc:
            warning = f"matplotlib unavailable for chart generation: {exc}"
            LOGGER.warning(warning)
            return None, {"type": "chart", "identifier": chart_identifier, "status": "skipped"}, warning

        metric = self._resolve_metric(chart_identifier)
        periods = self._resolve_periods(chart_identifier)
        analysis = self.analytics.analyze_metric_trend(
            ticker,
            metric,
            years_history=periods.get("history", 5),
            years_forecast=periods.get("forecast", 3),
        )

        if not analysis:
            warning = f"No historical data for {ticker} {metric}"
            LOGGER.warning(warning)
            return None, {
                "type": "chart",
                "identifier": chart_identifier,
                "status": "missing_data",
                "details": warning,
            }, warning

        forecasts = analysis.forecasts
        try:
            fig, ax = plt.subplots(figsize=(6, 3.5))
            historical_years = [year for year, _ in analysis.historical_data]
            historical_values = [value for _, value in analysis.historical_data]
            ax.plot(historical_years, historical_values, marker="o", label="Historical")

            if forecasts:
                forecast_years = [f.fiscal_year for f in forecasts]
                forecast_values = [f.predicted_value for f in forecasts]
                ax.plot(forecast_years, forecast_values, linestyle="--", marker="x", label="Forecast")
                for idx, f in enumerate(forecasts):
                    if f.confidence_interval_low is not None and f.confidence_interval_high is not None:
                        ax.fill_between(
                            [forecast_years[idx]],
                            [f.confidence_interval_low],
                            [f.confidence_interval_high],
                            color="gray",
                            alpha=0.2,
                        )

            ax.set_title(f"{ticker.upper()} {metric.replace('_', ' ').title()}")
            ax.set_xlabel("Year")
            ax.set_ylabel("Value")
            ax.legend()
            ax.grid(True, linestyle="--", alpha=0.3)

            tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            fig.tight_layout()
            fig.savefig(tmp_file.name, dpi=150)
            plt.close(fig)
            chart_path = tmp_file.name
        except Exception as exc:
            warning = f"Chart rendering failed: {exc}"
            LOGGER.error(warning, exc_info=True)
            return None, {
                "type": "chart",
                "identifier": chart_identifier,
                "status": "failed",
                "details": warning,
            }, warning

        audit_entry = {
            "type": "chart",
            "identifier": chart_identifier,
            "status": "generated",
            "metric": metric,
            "history_points": len(analysis.historical_data),
            "forecast_points": len(forecasts),
            "image_path": chart_path,
        }
        return chart_path, audit_entry, None

    @staticmethod
    def _resolve_metric(identifier: str) -> str:
        if "|" in identifier:
            identifier = identifier.split("|", 1)[0]
        alias = identifier.replace("Chart", "").replace("chart", "").strip()
        if ":" in alias:
            _, alias = alias.split(":", 1)
        mapping = {
            "revenuetrend": "revenue",
            "netincometrend": "net_income",
            "fcftrend": "free_cash_flow",
        }
        key = alias.lower().replace(" ", "").replace("-", "")
        return mapping.get(key, key or "revenue")

    @staticmethod
    def _resolve_periods(identifier: str) -> Dict[str, int]:
        periods = {"history": 5, "forecast": 3}
        if "|" not in identifier:
            return periods
        _, params_str = identifier.split("|", 1)
        params = params_str.split("|")
        for param in params:
            param_lower = param.lower()
            if param_lower.startswith("last_"):
                try:
                    periods["history"] = int(param_lower.split("_")[1])
                except Exception:
                    continue
            if param_lower.endswith("_years"):
                try:
                    periods["forecast"] = int(param_lower.split("_")[0])
                except Exception:
                    continue
        return periods

class TemplateProcessor:
    """Processes and manages report templates."""
    
    def __init__(self, db_path: Path):
        self.db_path = db_path
        self.extractor = TemplateExtractor()
        self.filler = TemplateFiller()
        self.custom_kpi_calculator = CustomKPICalculator(db_path)
        self.model_builder = ModelBuilder(db_path)
        self.chart_generator = ChartGenerator(db_path)
        self.source_tracer = SourceTracer(db_path)

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _lookup_company_name(self, ticker: str) -> Optional[str]:
        try:
            with sqlite3.connect(self.db_path) as conn:
                conn.row_factory = sqlite3.Row
                row = conn.execute(
                    """
                    SELECT company_name
                    FROM ticker_aliases
                    WHERE ticker = ?
                    """,
                    (ticker.upper(),),
                ).fetchone()
                if row:
                    return row["company_name"]
        except Exception as exc:
            LOGGER.debug("Company name lookup failed: %s", exc)
        return None

    def _build_context(self, ticker: str, user_id: str, context: Optional[Dict[str, Any]]) -> Dict[str, Any]:
        combined: Dict[str, Any] = {}
        combined.update(context or {})
        combined.setdefault("ticker", ticker.upper())
        combined.setdefault("user_id", user_id)
        combined.setdefault("company_name", self._lookup_company_name(ticker) or ticker.upper())
        combined.setdefault("generated_at", datetime.now(timezone.utc).isoformat())
        return combined

    def _resolve_placeholder(
        self,
        placeholder: Dict[str, Any],
        context: Dict[str, Any],
    ) -> Tuple[Any, Optional[Dict[str, Any]], Optional[str]]:
        data_type = placeholder.get("data_type", "text")
        identifier = placeholder.get("identifier") or placeholder.get("name", "")

        if data_type == "kpi":
            return self._resolve_kpi_placeholder(identifier, context)
        if data_type == "chart":
            return self._resolve_chart_placeholder(identifier, context)
        if data_type == "model":
            return self._resolve_model_placeholder(identifier, context)
        return self._resolve_text_placeholder(identifier, context)

    def _resolve_text_placeholder(
        self,
        identifier: str,
        context: Dict[str, Any],
    ) -> Tuple[Any, Optional[Dict[str, Any]], Optional[str]]:
        key = identifier.strip()
        value = context.get(key)
        if value is None:
            key_lower = key.lower()
            if key_lower in {"ticker"}:
                value = context.get("ticker")
            elif key_lower in {"company", "company_name"}:
                value = context.get("company_name")
            elif key_lower in {"date", "today"}:
                value = datetime.now().strftime("%Y-%m-%d")
            elif key_lower in context:
                value = context[key_lower]

        if value is None:
            warning = f"Placeholder '{identifier}' has no resolved value."
            audit = {"type": "text", "identifier": identifier, "status": "missing"}
            return f"[{identifier}]", audit, warning

        audit = {"type": "text", "identifier": identifier, "status": "resolved"}
        return value, audit, None

    def _resolve_kpi_placeholder(
        self,
        identifier: str,
        context: Dict[str, Any],
    ) -> Tuple[Any, Optional[Dict[str, Any]], Optional[str]]:
        ticker = context.get("ticker")
        user_id = context.get("user_id", "default")
        if not ticker:
            warning = f"No ticker provided for KPI placeholder '{identifier}'."
            return "[KPI unavailable]", {"type": "kpi", "identifier": identifier, "status": "missing_ticker"}, warning

        kpi = self.custom_kpi_calculator.get_kpi_by_name(user_id, identifier)
        if not kpi:
            warning = f"KPI '{identifier}' not found for user {user_id}."
            return "[KPI missing]", {"type": "kpi", "identifier": identifier, "status": "not_found"}, warning

        try:
            result = self.custom_kpi_calculator.calculate_kpi(kpi.kpi_id, ticker)
        except Exception as exc:
            warning = f"KPI calculation failed: {exc}"
            LOGGER.error(warning, exc_info=True)
            return "[KPI error]", {"type": "kpi", "identifier": identifier, "status": "error", "details": warning}, warning

        if result.error:
            warning = result.error
            audit = {
                "type": "kpi",
                "identifier": identifier,
                "status": "error",
                "details": warning,
            }
            return "[KPI error]", audit, warning

        value = result.formatted_value or (f"{result.value:,.2f}" if result.value is not None else "N/A")
        audit = {
            "type": "kpi",
            "identifier": identifier,
            "status": "resolved",
            "value": result.value,
            "formatted_value": value,
            "sources": result.sources,
            "metadata": result.metadata,
        }
        return value, audit, None

    def _resolve_chart_placeholder(
        self,
        identifier: str,
        context: Dict[str, Any],
    ) -> Tuple[Any, Optional[Dict[str, Any]], Optional[str]]:
        ticker = context.get("ticker")
        if not ticker:
            warning = f"No ticker provided for chart placeholder '{identifier}'."
            return "[Chart unavailable]", {"type": "chart", "identifier": identifier, "status": "missing_ticker"}, warning

        file_path, audit_entry, warning = self.chart_generator.generate(ticker, identifier, context)
        if not file_path:
            return "[Chart unavailable]", audit_entry, warning

        value = {
            "image_path": file_path,
            "text": f"See chart for {ticker}",
        }
        return value, audit_entry, warning

    def _resolve_model_placeholder(
        self,
        identifier: str,
        context: Dict[str, Any],
    ) -> Tuple[Any, Optional[Dict[str, Any]], Optional[str]]:
        user_id = context.get("user_id", "default")
        ticker = context.get("ticker")
        if not ticker:
            warning = f"No ticker provided for model placeholder '{identifier}'."
            return "[Model unavailable]", {"type": "model", "identifier": identifier, "status": "missing_ticker"}, warning

        parts = [part.strip() for part in identifier.split("|") if part.strip()]
        model_name = parts[0] if parts else identifier.strip()
        metric = parts[1] if len(parts) > 1 else None

        models = self.model_builder.list_models(user_id)
        model_match = next((m for m in models if m.name.lower() == model_name.lower()), None)
        if not model_match and models:
            # Attempt partial match
            model_match = next((m for m in models if model_name.lower() in m.name.lower()), None)

        if not model_match:
            warning = f"Model '{model_name}' not found for user {user_id}."
            return "[Model missing]", {"type": "model", "identifier": identifier, "status": "not_found"}, warning

        target_metric = metric or model_match.target_metric or (model_match.metrics[0] if model_match.metrics else None)
        if not target_metric:
            warning = f"No metric specified for model placeholder '{identifier}'."
            return "[Model metric missing]", {"type": "model", "identifier": identifier, "status": "missing_metric"}, warning

        try:
            explanation = self.model_builder.explain_model(model_match.model_id)
            latest_run = explanation.get("latest_run")
            if not latest_run or latest_run.get("ticker") != ticker:
                model_run = self.model_builder.run_forecast(
                    model_match.model_id,
                    ticker,
                    target_metric,
                    forecast_years=model_match.forecast_horizon or 3,
                )
                latest_run = model_run.to_dict()
        except Exception as exc:
            warning = f"Model execution failed: {exc}"
            LOGGER.error(warning, exc_info=True)
            return "[Model error]", {"type": "model", "identifier": identifier, "status": "error", "details": warning}, warning

        forecasts = (latest_run or {}).get("results", {}).get("forecasts", [])
        if not forecasts:
            warning = f"No forecasts produced for model '{model_name}'."
            return "[Model forecast missing]", {"type": "model", "identifier": identifier, "status": "no_forecast"}, warning

        latest_forecast = forecasts[-1]
        value = f"{model_name}: {target_metric} {latest_forecast.get('fiscal_year')} => {latest_forecast.get('predicted_value', 'N/A')}"
        audit = {
            "type": "model",
            "identifier": identifier,
            "status": "resolved",
            "model_id": model_match.model_id,
            "forecast_year": latest_forecast.get("fiscal_year"),
            "forecast_value": latest_forecast.get("predicted_value"),
            "driver_explanations": latest_run.get("driver_explanations", []),
        }
        return value, audit, None

    def _collect_placeholder_data(
        self,
        template: ReportTemplate,
        context: Dict[str, Any],
    ) -> Tuple[Dict[str, Any], List[Dict[str, Any]], List[str]]:
        data_map: Dict[str, Any] = {}
        audit_log: List[Dict[str, Any]] = []
        warnings: List[str] = []

        for placeholder in template.structure.get("placeholders", []):
            value, audit_entry, warning = self._resolve_placeholder(placeholder, context)
            if value is not None:
                data_map[placeholder.get("name", "")] = value
            if audit_entry:
                audit_log.append(audit_entry)
            if warning:
                warnings.append(warning)

        return data_map, audit_log, warnings
    
    def upload_template(
        self,
        user_id: str,
        name: str,
        file_path: Path,
        file_type: str
    ) -> ReportTemplate:
        """Upload and process a template."""
        # Extract structure
        structure = self.extractor.extract_structure(file_path, file_type)
        
        template_id = str(uuid.uuid4())
        now = datetime.now(timezone.utc)
        
        template = ReportTemplate(
            template_id=template_id,
            user_id=user_id,
            name=name,
            file_type=file_type,
            file_path=str(file_path),
            structure=structure,
            created_at=now,
        )
        
        # Save to database
        with sqlite3.connect(self.db_path) as conn:
            conn.execute(
                """
                INSERT INTO report_templates (template_id, user_id, name, file_type, file_path, structure, created_at)
                VALUES (?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    template.template_id,
                    template.user_id,
                    template.name,
                    template.file_type,
                    template.file_path,
                    json.dumps(template.structure),
                    template.created_at.isoformat(),
                ),
            )
            
            # Save placeholders
            for placeholder in structure.get("placeholders", []):
                placeholder_id = str(uuid.uuid4())
                conn.execute(
                    """
                    INSERT INTO template_placeholders (placeholder_id, template_id, placeholder_name, data_type, source_metric)
                    VALUES (?, ?, ?, ?, ?)
                    """,
                    (
                        placeholder_id,
                        template_id,
                        placeholder.get("name", ""),
                        placeholder.get("data_type", placeholder.get("type", "text")),
                        placeholder.get("identifier"),
                    ),
                )
            
            conn.commit()
        
        return template
    
    def get_template(self, template_id: str) -> Optional[ReportTemplate]:
        """Get a template by ID."""
        with sqlite3.connect(self.db_path) as conn:
            conn.row_factory = sqlite3.Row
            row = conn.execute(
                """
                SELECT template_id, user_id, name, file_type, file_path, structure, created_at
                FROM report_templates
                WHERE template_id = ?
                """,
                (template_id,),
            ).fetchone()
            
            if not row:
                return None
            
            return ReportTemplate(
                template_id=row["template_id"],
                user_id=row["user_id"],
                name=row["name"],
                file_type=row["file_type"],
                file_path=row["file_path"],
                structure=json.loads(row["structure"]) if row["structure"] else {},
                created_at=datetime.fromisoformat(row["created_at"]),
            )
    
    def list_templates(self, user_id: str) -> List[ReportTemplate]:
        """List all templates for a user."""
        with sqlite3.connect(self.db_path) as conn:
            conn.row_factory = sqlite3.Row
            rows = conn.execute(
                """
                SELECT template_id, user_id, name, file_type, file_path, structure, created_at
                FROM report_templates
                WHERE user_id = ?
                ORDER BY created_at DESC
                """,
                (user_id,),
            ).fetchall()
            
            return [
                ReportTemplate(
                    template_id=row["template_id"],
                    user_id=row["user_id"],
                    name=row["name"],
                    file_type=row["file_type"],
                    file_path=row["file_path"],
                    structure=json.loads(row["structure"]) if row["structure"] else {},
                    created_at=datetime.fromisoformat(row["created_at"]),
                )
                for row in rows
            ]
    
    def list_render_jobs(self, template_id: str, limit: int = 20) -> List[Dict[str, Any]]:
        """List recent render jobs for a template."""
        with sqlite3.connect(self.db_path) as conn:
            conn.row_factory = sqlite3.Row
            rows = conn.execute(
                """
                SELECT job_id, user_id, ticker, context, audit_log, output_path, created_at
                FROM template_render_jobs
                WHERE template_id = ?
                ORDER BY created_at DESC
                LIMIT ?
                """,
                (template_id, limit),
            ).fetchall()

            history: List[Dict[str, Any]] = []
            for row in rows:
                history.append(
                    {
                        "job_id": row["job_id"],
                        "user_id": row["user_id"],
                        "ticker": row["ticker"],
                        "context": json.loads(row["context"]) if row["context"] else {},
                        "audit_log": json.loads(row["audit_log"]) if row["audit_log"] else [],
                        "output_path": row["output_path"],
                        "created_at": row["created_at"],
                    }
                )
            return history

    def generate_from_template(
        self,
        template_id: str,
        ticker: str,
        output_path: Path,
        context: Optional[Dict[str, Any]] = None
    ) -> TemplateRenderResult:
        """Generate a report from a template with new data."""
        template = self.get_template(template_id)
        if not template:
            raise ValueError(f"Template {template_id} not found")

        user_id = (context or {}).get("user_id", "default")
        runtime_context = self._build_context(ticker, user_id, context)
        data_map, audit_log, warnings = self._collect_placeholder_data(template, runtime_context)

        # Ensure base context keys are available
        data_map.setdefault("ticker", runtime_context.get("ticker"))
        data_map.setdefault("company", runtime_context.get("company_name"))
        data_map.setdefault("generated_at", runtime_context.get("generated_at"))

        template_path = Path(template.file_path)
        success = self.filler.fill_template(template_path, template.file_type, data_map, output_path)

        job_id = str(uuid.uuid4())
        try:
            with sqlite3.connect(self.db_path) as conn:
                conn.execute(
                    """
                    INSERT INTO template_render_jobs (
                        job_id,
                        template_id,
                        user_id,
                        ticker,
                        context,
                        audit_log,
                        output_path,
                        created_at
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        job_id,
                        template_id,
                        user_id,
                        runtime_context.get("ticker"),
                        _json_dumps(runtime_context),
                        _json_dumps(audit_log),
                        str(output_path),
                        datetime.now(timezone.utc).isoformat(),
                    ),
                )
                conn.commit()
        except Exception as exc:
            LOGGER.error("Failed to record template render job: %s", exc, exc_info=True)

        return TemplateRenderResult(
            success=success,
            output_path=output_path,
            filled_data=data_map,
            audit_log=audit_log,
            warnings=warnings,
            job_id=job_id,
        )

